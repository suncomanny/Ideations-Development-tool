from __future__ import annotations

import json
import re
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from openpyxl import load_workbook

from .categories import Category
from .ideation_template import latest_gap_workbook
from .paths import ProjectPaths
from .research_tool import latest_prd_ideation_workbook
from .utils import newest_file, slugify, timestamp


SUMMARY_TABLE_LIMIT = 10


def _deck_spec_source_dir() -> Path:
    return Path(__file__).with_name("gate0_deck_specs")


def _copy_deck_specs(package_dir: Path) -> Path | None:
    source = _deck_spec_source_dir()
    if not source.exists():
        return None
    target = package_dir / "markdown_specs"
    if target.exists():
        shutil.rmtree(target)
    shutil.copytree(source, target)
    return target


@dataclass(frozen=True)
class Gate0DeckResult:
    deck_path: Path
    package_dir: Path
    config_path: Path
    slide_data_path: Path
    audit_path: Path
    spec_snapshot_dir: Path | None
    warnings: list[str]


def latest_research_report(paths: ProjectPaths, category: Category) -> Path | None:
    patterns = [
        f"{category.slug}_*_completed_rows*.xlsx",
        f"{category.slug}_*completed*.xlsx",
    ]
    return newest_file(paths.research_report_category_outputs(category.slug), patterns) or newest_file(
        paths.research_report_outputs,
        patterns,
    )


def _clean_header(value: Any) -> str:
    text = str(value or "").strip().replace("*", "")
    text = re.sub(r"\s+", " ", text)
    return text


def _norm_key(value: Any) -> str:
    text = _clean_header(value).lower()
    text = text.replace("/", " ")
    text = re.sub(r"[^a-z0-9]+", "_", text)
    return text.strip("_")


def _clean_text(value: Any, max_chars: int | None = None) -> str:
    if value is None:
        return ""
    text = str(value).replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text).strip()
    if max_chars and len(text) > max_chars:
        return text[: max_chars - 3].rstrip() + "..."
    return text


def _display_slide_title(value: Any) -> str:
    text = _clean_text(value)
    if ":" in text:
        prefix, remainder = text.split(":", 1)
        if any(token in prefix.lower() for token in ["candidate", "opportunity", "winner"]):
            text = remainder.strip()
    parts = [_clean_text(part) for part in text.split(";") if _clean_text(part)]
    if len(parts) >= 3:
        text = " - ".join(parts[:4])
    return _clean_text(text, 68)


def _as_float(value: Any) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).replace("$", "").replace(",", "").replace("%", "").strip()
    if not text:
        return None
    try:
        return float(text)
    except ValueError:
        return None


def _fmt_money(value: Any) -> str:
    number = _as_float(value)
    if number is None:
        return "N/A"
    if abs(number) >= 1_000_000:
        return f"${number / 1_000_000:.1f}M"
    if abs(number) >= 1_000:
        return f"${number / 1_000:.0f}K"
    return f"${number:,.0f}"


def _fmt_value(value: Any) -> str:
    if value in (None, ""):
        return "N/A"
    if isinstance(value, (int, float)):
        return f"{value:,.0f}" if float(value).is_integer() else f"{value:,.2f}"
    return _clean_text(value, 80)


def _read_table(path: Path, sheet_name: str, header_row: int = 1, start_row: int | None = None) -> list[dict[str, Any]]:
    if not path or not path.exists():
        return []
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        if sheet_name not in workbook.sheetnames:
            return []
        ws = workbook[sheet_name]
        headers = [_clean_header(ws.cell(header_row, col).value) for col in range(1, ws.max_column + 1)]
        rows: list[dict[str, Any]] = []
        first_row = start_row or header_row + 1
        for row_index in range(first_row, ws.max_row + 1):
            values = [ws.cell(row_index, col).value for col in range(1, len(headers) + 1)]
            if not any(value not in (None, "") for value in values):
                continue
            record = {
                headers[col_index]: values[col_index]
                for col_index in range(len(headers))
                if headers[col_index]
            }
            rows.append(record)
        return rows
    finally:
        workbook.close()


def _read_prd_rows(path: Path | None) -> list[dict[str, Any]]:
    if not path or not path.exists():
        return []
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        if "Ideations" not in workbook.sheetnames:
            return []
        ws = workbook["Ideations"]
        headers = [_clean_header(ws.cell(3, col).value) for col in range(1, ws.max_column + 1)]
        rows: list[dict[str, Any]] = []
        for row_index in range(4, ws.max_row + 1):
            values = [ws.cell(row_index, col).value for col in range(1, len(headers) + 1)]
            if not any(value not in (None, "") for value in values):
                continue
            record = {
                _norm_key(headers[col_index]): values[col_index]
                for col_index in range(len(headers))
                if headers[col_index]
            }
            if record.get("ideation_name"):
                rows.append(record)
        return rows
    finally:
        workbook.close()


def _read_source_mapping(path: Path | None) -> dict[str, dict[str, Any]]:
    rows = _read_table(path, "Source Mapping")
    mapping: dict[str, dict[str, Any]] = {}
    for row in rows:
        name = _clean_text(row.get("Ideation Name"))
        if name:
            mapping[name.lower()] = row
    return mapping


def _read_research_summary(path: Path | None) -> dict[str, dict[str, Any]]:
    if not path or not path.exists():
        return {}
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        if "Summary" not in workbook.sheetnames:
            return {}
        ws = workbook["Summary"]
        header_row = None
        for row_index in range(1, ws.max_row + 1):
            if _clean_header(ws.cell(row_index, 1).value).lower() == "row" and _clean_header(ws.cell(row_index, 2).value).lower() == "ideation":
                header_row = row_index
                break
        if not header_row:
            return {}
        headers = [_clean_header(ws.cell(header_row, col).value) for col in range(1, ws.max_column + 1)]
        summary: dict[str, dict[str, Any]] = {}
        for row_index in range(header_row + 1, ws.max_row + 1):
            values = [ws.cell(row_index, col).value for col in range(1, len(headers) + 1)]
            if not any(value not in (None, "") for value in values):
                continue
            record = {
                _norm_key(headers[col_index]): values[col_index]
                for col_index in range(len(headers))
                if headers[col_index]
            }
            name = _clean_text(record.get("ideation"))
            if name:
                summary[name.lower()] = record
        return summary
    finally:
        workbook.close()


def _find_row(ws, label: str) -> int | None:
    target = label.strip().lower()
    for row_index in range(1, ws.max_row + 1):
        if _clean_header(ws.cell(row_index, 1).value).lower() == target:
            return row_index
    return None


def _extract_key_pairs(ws, start_row: int, end_row: int) -> dict[str, Any]:
    values: dict[str, Any] = {}
    for row_index in range(start_row, min(end_row, ws.max_row) + 1):
        pairs = [(1, 2), (3, 4), (5, 6)]
        for key_col, value_col in pairs:
            key = _clean_header(ws.cell(row_index, key_col).value)
            if not key:
                continue
            values[_norm_key(key)] = ws.cell(row_index, value_col).value
    return values


def _extract_table_after(ws, title: str) -> list[dict[str, Any]]:
    title_row = _find_row(ws, title)
    if not title_row:
        return []
    header_row = None
    for row_index in range(title_row + 1, min(ws.max_row, title_row + 5) + 1):
        first = _clean_header(ws.cell(row_index, 1).value)
        second = _clean_header(ws.cell(row_index, 2).value)
        if first and second and not first.lower().startswith("section "):
            header_row = row_index
            break
    if not header_row:
        return []
    headers = [_clean_header(ws.cell(header_row, col).value) for col in range(1, ws.max_column + 1)]
    rows: list[dict[str, Any]] = []
    for row_index in range(header_row + 1, ws.max_row + 1):
        first = _clean_header(ws.cell(row_index, 1).value)
        if not first:
            break
        if first.lower().startswith("section ") or first in {
            "Numeric Target Positioning",
            "Category Optimization Summary",
            "Optimization Profile",
            "Primary Decision Drivers",
            "Secondary Decision Drivers",
            "Validate Before Over-Weighting",
            "Highest-Impact Vendor Requests",
            "Recommendations",
            "Notes",
        }:
            break
        if first == "No rows.":
            break
        values = [ws.cell(row_index, col).value for col in range(1, len(headers) + 1)]
        record = {
            headers[col_index]: values[col_index]
            for col_index in range(len(headers))
            if headers[col_index]
        }
        rows.append(record)
    return rows


def _extract_recommendation_bullets(ws) -> list[str]:
    start = _find_row(ws, "Recommendations")
    if not start:
        return []
    bullets: list[str] = []
    for row_index in range(start + 1, min(ws.max_row, start + 10) + 1):
        text = _clean_text(ws.cell(row_index, 1).value)
        if not text:
            break
        if text == "Notes":
            break
        bullets.append(text.lstrip("- ").strip())
    return bullets


def _read_research_details(path: Path | None) -> dict[str, dict[str, Any]]:
    if not path or not path.exists():
        return {}
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        details: dict[str, dict[str, Any]] = {}
        for ws in workbook.worksheets:
            if ws.title == "Summary":
                continue
            name = _clean_text(ws.cell(1, 1).value)
            if not name:
                continue
            record = {
                "overall_read": ws.cell(5, 2).value if _clean_header(ws.cell(5, 1).value) == "Overall Read" else None,
                "why_not_stronger": ws.cell(6, 2).value if _clean_header(ws.cell(6, 1).value) == "Why Not Stronger" else None,
                "anchor_context": _extract_key_pairs(ws, 12, 29),
                "pricing": _extract_key_pairs(ws, 63, 76),
                "optimization": _extract_key_pairs(ws, 119, 123),
                "gate_readiness_summary": ws.cell(34, 2).value if _clean_header(ws.cell(34, 1).value) == "Gate Readiness Summary" else None,
                "category_optimization_summary": ws.cell(119, 2).value if _clean_header(ws.cell(119, 1).value) == "Category Optimization Summary" else None,
                "verified_amazon_competitors": _extract_table_after(ws, "Section B - Amazon Competitors (Verified Listings)"),
                "verified_direct_competitors": _extract_table_after(ws, "Section C - Brick-and-Mortar / Direct Competitors (Verified Listings)"),
                "inferred_competitors": _extract_table_after(ws, "Section D - Inferred Competitors / Needs Verification"),
                "feature_signals": _extract_table_after(ws, "Section F - Feature / Certification Signals"),
                "primary_decision_drivers": _extract_table_after(ws, "Primary Decision Drivers"),
                "vendor_requests": _extract_table_after(ws, "Highest-Impact Vendor Requests"),
                "recommendations": _extract_recommendation_bullets(ws),
                "prd_prefill": _extract_key_pairs(ws, 173, 184),
            }
            details[name.lower()] = record
        return details
    finally:
        workbook.close()


def _top_line_review_rows(rows: list[dict[str, Any]], limit: int = 8) -> list[dict[str, Any]]:
    filtered = [
        row for row in rows
        if not str(row.get("Family Part Number") or "").lower().startswith("(no approved")
    ]
    return sorted(filtered, key=lambda row: _as_float(row.get("Total Revenue")) or 0, reverse=True)[:limit]


def _merge_ideations(
    *,
    prd_rows: list[dict[str, Any]],
    source_mapping: dict[str, dict[str, Any]],
    research_summary: dict[str, dict[str, Any]],
    research_details: dict[str, dict[str, Any]],
) -> list[dict[str, Any]]:
    merged: list[dict[str, Any]] = []
    for index, row in enumerate(prd_rows, start=1):
        name = _clean_text(row.get("ideation_name"))
        key = name.lower()
        source = source_mapping.get(key, {})
        summary = research_summary.get(key, {})
        details = research_details.get(key, {})
        prd_prefill = details.get("prd_prefill") or {}
        specs = {
            "voltage": row.get("voltage") or prd_prefill.get("voltage"),
            "wattage": row.get("wattage_primary") or prd_prefill.get("wattage_primary"),
            "cct": row.get("cct_primary") or prd_prefill.get("cct_primary"),
            "lumens": row.get("lumens_target") or prd_prefill.get("lumens_target"),
            "cri": row.get("cri") or prd_prefill.get("cri"),
            "dimming": row.get("dimming_type") or row.get("dimmable") or prd_prefill.get("dimming_type"),
            "form_factor": row.get("size_form_factor") or prd_prefill.get("size_form_factor"),
            "mounting": row.get("mounting_type") or prd_prefill.get("mounting_type"),
            "moisture_rating": row.get("moisture_rating") or prd_prefill.get("moisture_rating"),
            "certifications": row.get("certifications") or prd_prefill.get("certifications"),
            "target_msrp": row.get("target_msrp") or (details.get("pricing") or {}).get("target_msrp"),
        }
        merged.append({
            "slide_number": index,
            "name": name,
            "category": _clean_text(row.get("category")),
            "subcategory": _clean_text(row.get("subcategory")),
            "strategy": _clean_text(row.get("strategy")),
            "reference_sku": _clean_text(row.get("sunco_reference_sku")),
            "specs": {key: _clean_text(value) for key, value in specs.items() if _clean_text(value)},
            "source_why": _clean_text(source.get("Why selected"), 220),
            "source_evidence": _clean_text(source.get("Key evidence used"), 700),
            "research_summary": summary,
            "research_details": details,
        })
    return merged


def _best_competitor(item: dict[str, Any]) -> dict[str, Any]:
    details = item.get("research_details") or {}
    for source_key in ["verified_amazon_competitors", "verified_direct_competitors", "inferred_competitors"]:
        rows = details.get(source_key) or []
        if rows:
            row = rows[0]
            return {
                "brand": _clean_text(row.get("Brand")),
                "product": _clean_text(row.get("Product"), 140),
                "identifier": _clean_text(row.get("Identifier") or row.get("Product Code") or row.get("Model")),
                "channel": _clean_text(row.get("Channel") or row.get("Likely Channel")),
                "price": _clean_text(row.get("Price")),
            }
    return {}


def _feature_signal_bullets(item: dict[str, Any], limit: int = 4) -> list[str]:
    details = item.get("research_details") or {}
    rows = details.get("primary_decision_drivers") or details.get("feature_signals") or []
    bullets: list[str] = []
    for row in rows:
        label = row.get("Driver") or row.get("Label")
        signal = row.get("Signal")
        coverage = row.get("Coverage %")
        if label:
            suffix = f" ({signal}" if signal else ""
            if coverage not in (None, ""):
                suffix += f", {coverage}% coverage" if suffix else f" ({coverage}% coverage"
            suffix += ")" if suffix else ""
            bullets.append(f"{_clean_text(label, 48)}{suffix}")
        if len(bullets) >= limit:
            break
    return bullets


def _vendor_request_bullets(item: dict[str, Any], limit: int = 3) -> list[str]:
    rows = (item.get("research_details") or {}).get("vendor_requests") or []
    bullets: list[str] = []
    for row in rows:
        request = _clean_text(row.get("Request"), 120)
        reason = _clean_text(row.get("Reason"), 90)
        if request:
            bullets.append(f"{request} - {reason}" if reason else request)
        if len(bullets) >= limit:
            break
    if bullets:
        return bullets
    recommendations = (item.get("research_details") or {}).get("recommendations") or []
    return [_clean_text(text, 140) for text in recommendations[:limit]]


def _summary_rows(ideations: list[dict[str, Any]]) -> list[list[str]]:
    rows: list[list[str]] = []
    for item in ideations[:SUMMARY_TABLE_LIMIT]:
        summary = item.get("research_summary") or {}
        specs = item.get("specs") or {}
        key_spec = specs.get("form_factor") or specs.get("wattage") or specs.get("cct") or "See detail slide"
        next_step = "PRD/RFQ review" if "favor" in str(summary.get("outlook") or "").lower() else "Evidence review"
        rows.append([
            _clean_text(item.get("name"), 46),
            _clean_text(item.get("strategy"), 20),
            _clean_text(key_spec, 34),
            _clean_text(summary.get("outlook") or "needs review", 18),
            _clean_text(summary.get("confidence") or "n/a", 14),
            next_step,
        ])
    return rows


def _build_config(
    *,
    category: Category,
    run_timestamp: str,
    ideations: list[dict[str, Any]],
    line_review_rows: list[dict[str, Any]],
    paths_used: dict[str, str],
) -> dict[str, Any]:
    total_revenue = sum((_as_float(row.get("Total Revenue")) or 0) for row in line_review_rows)
    favorable_count = sum(1 for item in ideations if "favor" in str((item.get("research_summary") or {}).get("outlook") or "").lower())
    config = {
        "meta": {
            "author": "Sunco Product Opportunity Engine",
            "title": f"{category.run_name} Gate 0 Opportunity Review",
            "outputFileName": f"{category.slug}_gate0_deck_{run_timestamp}.pptx",
            "slideDataFile": f"{category.slug}_gate0_slide_data_{run_timestamp}.json",
            "sourceMode": "Step 1/2/3 outputs only",
        },
        "marketSnapshot": {
            "title": f"{category.run_name} Gate 0 Opportunity Snapshot",
            "subtitle": "Leadership opportunity review generated from the current 1/2/3 workflow",
            "marketOverview": {
                "header": "Evidence Base",
                "bullets": [
                    f"Step 1 gap workbook: {Path(paths_used.get('step1') or '').name}",
                    f"Step 2 SKU-level ideations: {len(ideations)} row(s)",
                    f"Step 3 analyzed rows found: {sum(1 for item in ideations if item.get('research_summary'))}",
                    f"Current line-review rows included: {len(line_review_rows)}",
                ],
            },
            "barChart": {
                "title": "Workflow Evidence Counts",
                "labels": ["Step 2 Rows", "Step 3 Matched", "Line Rows"],
                "values": [
                    len(ideations),
                    sum(1 for item in ideations if item.get("research_summary")),
                    len(line_review_rows),
                ],
            },
            "pieChart": {
                "title": "Gate 0 Row Mix",
                "labels": ["Favorable", "Needs Review"],
                "values": [favorable_count, max(0, len(ideations) - favorable_count)],
                "colors": ["00A5B5", "94A3B8"],
            },
            "segmentPriorities": {
                "header": "What Leadership Should Decide",
                "items": [
                    {"label": "Pursue: ", "desc": "Approve selected opportunities to move into PRD/RFQ work."},
                    {"label": "Hold: ", "desc": "Pause rows with weak evidence, missing pricing, or poor line fit."},
                    {"label": "Research Gap: ", "desc": "Call out missing Stackline, pricing, vendor, or competitor evidence before PRD scope."},
                ],
            },
            "topOpportunities": {
                "header": "Top Opportunities",
                "bullets": [_clean_text(item.get("name"), 90) for item in ideations[:5]] or ["No Step 2 ideations found."],
            },
            "demandSignals": {
                "header": "Demand / Fit Signals",
                "bullets": [
                    f"{favorable_count} ideation(s) currently read favorable in Step 3." if favorable_count else "Step 3 evidence should be reviewed row by row before approval.",
                    f"Existing line revenue represented in deck: {_fmt_money(total_revenue)}.",
                    "Feature and certification drivers come from Step 3 Section F, not from the legacy deck data sources.",
                ],
            },
            "recommendedActions": {
                "header": "Recommended Gate 0 Action",
                "bullets": [
                    "Use this deck to align on which opportunities deserve PRD/RFQ work.",
                    "Do not treat this as PO approval; vendor quotes, samples, and pack-size work happen later.",
                ],
            },
            "statCallouts": [
                {"value": str(len(ideations)), "label": "SKU-Level Ideations"},
                {"value": str(len(line_review_rows)), "label": "Current Line Rows"},
                {"value": _fmt_money(total_revenue), "label": "Current Line Revenue"},
            ],
            "sourceFootnote": "Sources: current Step 1 gap workbook, Step 2 ideation workbook, and Step 3 research report only.",
        },
        "summarySlide": {
            "title": f"{category.run_name} Opportunity Line Review",
            "stats": [
                {"label": "Step 2 Rows", "value": len(ideations)},
                {"label": "Step 3 Matched", "value": sum(1 for item in ideations if item.get("research_summary"))},
                {"label": "Current Line Rows", "value": len(line_review_rows)},
                {"label": "Favorable Reads", "value": favorable_count},
            ],
            "tableTitle": "Opportunity Shortlist",
            "tableHeaders": ["Ideation", "Strategy", "Key Spec", "Outlook", "Confidence", "Next Step"],
            "ideations": _summary_rows(ideations),
        },
        "productSlides": {
            "titles": {str(item["slide_number"]): item["name"] for item in ideations},
            "strategies": {},
            "strategyNotes": {},
        },
        "sourceFiles": paths_used,
    }
    for item in ideations:
        number = str(item["slide_number"])
        details = item.get("research_details") or {}
        notes = [
            f"Source why: {item.get('source_why') or 'Not available'}",
            f"Evidence: {item.get('source_evidence') or 'Not available'}",
            f"Overall read: {_clean_text(details.get('overall_read')) or 'Not available'}",
            f"Gate readiness: {_clean_text(details.get('gate_readiness_summary')) or 'Not available'}",
        ]
        config["productSlides"]["strategies"][number] = _clean_text(details.get("overall_read") or item.get("source_why") or item.get("strategy"), 520)
        config["productSlides"]["strategyNotes"][number] = "\n".join(f"- {note}" for note in notes if note)
    return config


def _build_slide_data(ideations: list[dict[str, Any]]) -> list[dict[str, Any]]:
    slide_data: list[dict[str, Any]] = []
    for item in ideations:
        specs = item.get("specs") or {}
        summary = item.get("research_summary") or {}
        details = item.get("research_details") or {}
        pricing = details.get("pricing") or {}
        target_price = _as_float(pricing.get("target_msrp") or specs.get("target_msrp"))
        competitor = _best_competitor(item)
        identifier = competitor.get("identifier") or ""
        asin_match = re.search(r"\bB0[A-Z0-9]{8}\b|\bB[A-Z0-9]{9}\b", identifier)
        asin = asin_match.group(0) if asin_match else identifier.replace("ASIN", "").strip()
        competitor_for_legacy = {
            "asin": asin or "TBD",
            "url": "",
            "brand": competitor.get("brand") or "TBD",
            "price": competitor.get("price") or None,
            "parentRevenue": 0,
            "parentSales": 0,
            "rating": None,
            "reviews": None,
            "segment": competitor.get("channel") or "Verified competitor",
            "dataSource": "step3_research_report",
            "tacos": None,
            "product": competitor.get("product") or "",
            "identifier": identifier,
            "channel": competitor.get("channel") or "",
        }
        direct_competitor = {
            "segment": competitor.get("channel") or "Verified competitor",
            "productCode": identifier,
            "brand": competitor.get("brand") or "",
            "price": competitor.get("price") or None,
            "productName": competitor.get("product") or "",
            "name": competitor.get("product") or "",
            "url": "",
            "product": competitor.get("product") or "",
            "identifier": identifier,
            "channel": competitor.get("channel") or "",
        }
        slide_data.append({
            "slideNumber": item["slide_number"],
            "familyName": item["name"],
            "comparableSKU": item.get("reference_sku") or "TBD",
            "size": specs.get("form_factor") or "TBD",
            "wattages": specs.get("wattage") or "TBD",
            "ccts": specs.get("cct") or "TBD",
            "voltage": specs.get("voltage") or "TBD",
            "dimming": specs.get("dimming") or "TBD",
            "cri": specs.get("cri") or "TBD",
            "moistureRating": specs.get("moisture_rating") or "TBD",
            "certifications": specs.get("certifications") or "TBD",
            "performance": {
                "shopifyRevenue": 0,
                "amazonRevenue": 0,
                "salePrice": target_price,
                "marginDollars": None,
                "marginPercent": None,
                "unitsSold": 0,
                "orders": 0,
                "customers": 0,
                "shopifyL30d": 0,
                "imageSrc": None,
            },
            "variants": [{
                "skuName": item["name"],
                "potentialSKU": item["strategy"],
                "features": "; ".join(value for value in specs.values() if value),
                "differentiator": _feature_signal_bullets(item),
            }],
            "competitor": competitor_for_legacy,
            "directCompetitor": direct_competitor,
            "gate0": {
                "outlook": summary.get("outlook"),
                "confidence": summary.get("confidence"),
                "amazonG2": summary.get("amazon_g2"),
                "amazonEvidence": summary.get("amazon_evidence"),
                "overallRead": details.get("overall_read"),
                "whyNotStronger": details.get("why_not_stronger"),
                "categoryOptimizationSummary": details.get("category_optimization_summary"),
                "featureSignals": _feature_signal_bullets(item),
                "vendorRequests": _vendor_request_bullets(item),
                "sourceEvidence": item.get("source_evidence"),
            },
        })
    return slide_data


def _add_textbox(slide, left, top, width, height, text: str, *, font_size: int = 12, bold: bool = False, color=None, fill=None):
    from pptx.util import Pt

    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape


def _add_bullets(slide, left, top, width, height, bullets: list[str], *, font_size: int = 11):
    from pptx.util import Pt

    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = _clean_text(bullet, 180)
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.space_after = Pt(3)
    return shape


def _add_header(slide, title: str, subtitle: str | None = None):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RGBColor(0, 165, 181)
    _add_textbox(slide, Inches(0.45), Inches(0.18), Inches(12.2), Inches(0.35), title, font_size=20, bold=True, color=RGBColor(30, 41, 59))
    if subtitle:
        _add_textbox(slide, Inches(0.45), Inches(0.57), Inches(12.2), Inches(0.25), subtitle, font_size=9, color=RGBColor(100, 116, 139))


def _add_footer(slide, page_label: str):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    _add_textbox(slide, Inches(10.9), Inches(7.18), Inches(1.9), Inches(0.2), page_label, font_size=8, bold=True, color=RGBColor(0, 138, 151))


def _add_section_box(slide, left, top, width, height, title: str, body: str | list[str]):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 249, 250)
    box.line.color.rgb = RGBColor(226, 232, 240)
    _add_textbox(slide, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.22), title, font_size=10, bold=True, color=RGBColor(0, 138, 151))
    if isinstance(body, list):
        _add_bullets(slide, left + Inches(0.16), top + Inches(0.38), width - Inches(0.28), height - Inches(0.44), body, font_size=9)
    else:
        _add_textbox(slide, left + Inches(0.16), top + Inches(0.38), width - Inches(0.28), height - Inches(0.44), _clean_text(body, 700), font_size=9, color=RGBColor(51, 65, 85))


def _add_table(slide, left, top, width, height, headers: list[str], rows: list[list[str]], *, font_size: int = 8):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    row_count = max(1, len(rows)) + 1
    table_shape = slide.shapes.add_table(row_count, len(headers), left, top, width, height)
    table = table_shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 165, 181)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
    for row_index, row in enumerate(rows or [["No rows available"] + [""] * (len(headers) - 1)], start=1):
        for col_index in range(len(headers)):
            cell = table.cell(row_index, col_index)
            cell.text = _clean_text(row[col_index] if col_index < len(row) else "", 110)
            if row_index % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 250, 251)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor(30, 41, 59)
    return table_shape


def _write_pptx(
    *,
    output_path: Path,
    category: Category,
    config: dict[str, Any],
    slide_data: list[dict[str, Any]],
    line_review_rows: list[dict[str, Any]],
) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError(
            "Gate 0 deck generation requires the python-pptx package. "
            "Install it in the Python environment used by this tool with: python -m pip install python-pptx"
        ) from exc

    COLORS = {
        "teal": RGBColor(0, 165, 181),
        "teal_dark": RGBColor(0, 138, 151),
        "white": RGBColor(255, 255, 255),
        "off_white": RGBColor(247, 249, 250),
        "table_alt": RGBColor(240, 250, 251),
        "light_gray": RGBColor(232, 236, 239),
        "med_gray": RGBColor(148, 163, 184),
        "dark": RGBColor(30, 41, 59),
        "body": RGBColor(51, 65, 85),
        "navy": RGBColor(51, 65, 85),
    }
    SCALE = 4 / 3
    FONT_SCALE = 4 / 3

    def inch(value: float):
        return Inches(value * SCALE)

    def scaled_pt(value: float):
        from pptx.util import Pt

        return Pt(value * FONT_SCALE)

    def tx(slide, x, y, w, h, text: str, *, size=8, bold=False, italic=False, underline=False, color=None, align="left"):
        shape = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.TOP
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = paragraph.add_run()
        run.text = text or ""
        run.font.name = "Calibri"
        run.font.size = scaled_pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        run.font.color.rgb = color or COLORS["body"]
        return shape

    def paragraph_box(slide, x, y, w, h, lines: list[str], *, size=7.5, color=None):
        text = "\n".join(f"•     {_clean_text(line, 145)}" for line in lines if _clean_text(line))
        return tx(slide, x, y, w, h, text, size=size, color=color or COLORS["body"])

    def rect(slide, x, y, w, h, fill, line=None):
        shape = slide.shapes.add_shape(1, inch(x), inch(y), inch(w), inch(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is not None:
            shape.line.color.rgb = line
        else:
            shape.line.fill.background()
        return shape

    def accent(slide):
        rect(slide, 0, 0, 10, 0.05, COLORS["teal"])

    def footer(slide):
        rect(slide, 0, 5.35, 10, 0.275, COLORS["off_white"])
        tx(slide, 8.2, 5.39, 1.5, 0.16, "Sunco", size=10, bold=True, color=COLORS["teal_dark"], align="right")

    def section_header(slide, x, y, w, title: str, *, size=11):
        tx(slide, x, y, w, 0.18, title, size=size, bold=True, underline=True, color=COLORS["teal"])

    def write_table(slide, x, y, col_widths: list[float], row_height: float, headers: list[str], rows: list[list[str]], *, header_fill=None, font_size=6.5):
        row_count = max(1, len(rows)) + 1
        table_shape = slide.shapes.add_table(
            row_count,
            len(headers),
            inch(x),
            inch(y),
            inch(sum(col_widths)),
            inch(row_height * row_count),
        )
        table = table_shape.table
        for index, width in enumerate(col_widths):
            table.columns[index].width = inch(width)
        for row in table.rows:
            row.height = inch(row_height)
        for col_index, header in enumerate(headers):
            cell = table.cell(0, col_index)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill or COLORS["teal"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = scaled_pt(font_size)
                    run.font.bold = True
                    run.font.color.rgb = COLORS["white"]
        body_rows = rows or [["No rows."] + [""] * (len(headers) - 1)]
        for row_index, row in enumerate(body_rows, start=1):
            for col_index in range(len(headers)):
                cell = table.cell(row_index, col_index)
                cell.text = _clean_text(row[col_index] if col_index < len(row) else "", 95)
                if row_index % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS["table_alt"]
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT if col_index in {0, len(headers) - 1} else PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = "Calibri"
                        run.font.size = scaled_pt(font_size)
                        run.font.color.rgb = COLORS["dark"]
        return table_shape

    def source_footnote(slide, text: str | None = None):
        tx(
            slide,
            0.4,
            5.18,
            9.2,
            0.12,
            text or "Sources: current Step 1 gap workbook, Step 2 ideation workbook, and Step 3 research report.",
            size=5.5,
            italic=True,
            color=COLORS["med_gray"],
        )

    def add_snapshot_slide(prs):
        slide = prs.slides.add_slide(blank)
        ms = config["marketSnapshot"]
        accent(slide)
        tx(slide, 0.4, 0.1, 9.2, 0.3, ms["title"], size=18, bold=True, color=COLORS["dark"])
        tx(slide, 0.4, 0.42, 9.2, 0.18, ms["subtitle"], size=9, italic=True, color=COLORS["med_gray"])

        section_header(slide, 0.4, 0.75, 4.4, ms["marketOverview"]["header"])
        paragraph_box(slide, 0.4, 0.98, 4.5, 0.85, ms["marketOverview"]["bullets"], size=6.8)

        section_header(slide, 0.4, 1.85, 4.4, "Workflow Evidence Counts", size=8)
        chart_labels = ms.get("barChart", {}).get("labels") or []
        chart_values = ms.get("barChart", {}).get("values") or []
        max_value = max(chart_values or [1])
        for idx, (label, value) in enumerate(zip(chart_labels, chart_values)):
            y = 2.12 + idx * 0.26
            tx(slide, 0.42, y, 1.35, 0.14, label, size=6.2, color=COLORS["dark"], align="right")
            rect(slide, 1.9, y + 0.03, max(0.06, 2.55 * (float(value) / max_value)), 0.08, COLORS["teal"])
            tx(slide, 4.55, y, 0.35, 0.14, str(value), size=6.2, color=COLORS["dark"], align="right")

        section_header(slide, 5.2, 0.75, 4.3, ms["segmentPriorities"]["header"])
        priority_lines = [f"{item.get('label', '')}{item.get('desc', '')}" for item in ms["segmentPriorities"]["items"]]
        paragraph_box(slide, 5.2, 0.98, 4.3, 0.7, priority_lines, size=6.8)

        section_header(slide, 5.2, 1.78, 4.3, ms["topOpportunities"]["header"])
        paragraph_box(slide, 5.2, 2.0, 4.3, 0.86, ms["topOpportunities"]["bullets"], size=6.8)

        section_header(slide, 5.2, 2.95, 4.3, ms["demandSignals"]["header"])
        paragraph_box(slide, 5.2, 3.18, 4.3, 0.7, ms["demandSignals"]["bullets"], size=6.8)

        section_header(slide, 5.2, 4.02, 4.3, ms["recommendedActions"]["header"])
        paragraph_box(slide, 5.2, 4.25, 4.3, 0.45, ms["recommendedActions"]["bullets"], size=6.8)

        for index, stat in enumerate(ms["statCallouts"][:3]):
            x = 5.2 + index * 1.55
            rect(slide, x, 4.75, 1.42, 0.58, COLORS["off_white"], COLORS["light_gray"])
            tx(slide, x, 4.86, 1.42, 0.2, stat["value"], size=14, bold=True, color=COLORS["teal"], align="center")
            tx(slide, x + 0.08, 5.12, 1.26, 0.12, stat["label"], size=5.5, color=COLORS["med_gray"], align="center")
        source_footnote(slide, ms.get("sourceFootnote"))
        footer(slide)

    def add_line_review_slide(prs):
        slide = prs.slides.add_slide(blank)
        accent(slide)
        tx(slide, 0.4, 0.18, 9.2, 0.42, "Current Sunco Line Review", size=24, bold=True, color=COLORS["dark"], align="center")
        line_rows = [
            [
                _clean_text(row.get("Family Part Number"), 22),
                _clean_text(row.get("Product Title"), 38),
                _clean_text(row.get("Vendor"), 14),
                _fmt_money(row.get("Amazon Revenue")),
                _fmt_money(row.get("Shopify Revenue")),
                _clean_text(row.get("Pack Sizes Available"), 18),
            ]
            for row in line_review_rows[:14]
        ]
        write_table(
            slide,
            0.3,
            0.9,
            [1.15, 3.0, 1.0, 1.0, 1.0, 1.45],
            0.18,
            ["Family", "Product", "Vendor", "Amazon Rev", "Shopify Rev", "Pack Sizes"],
            line_rows,
            font_size=5.8,
        )
        source_footnote(slide)
        footer(slide)

    def add_summary_slide(prs):
        slide = prs.slides.add_slide(blank)
        ss = config["summarySlide"]
        accent(slide)
        tx(slide, 0.4, 0.18, 9.2, 0.42, ss["title"], size=24, bold=True, color=COLORS["dark"], align="center")
        box_x = 0.3
        for index, stat in enumerate(ss["stats"]):
            y = 0.8 + index * 0.42
            rect(slide, box_x, y, 0.06, 0.36, COLORS["teal"])
            rect(slide, box_x + 0.06, y, 3.14, 0.36, COLORS["off_white"], COLORS["light_gray"])
            tx(slide, box_x + 0.15, y + 0.1, 2.25, 0.12, stat["label"], size=9.5, bold=True, color=COLORS["teal"])
            tx(slide, box_x + 2.5, y + 0.06, 0.55, 0.16, str(stat["value"]), size=14, bold=True, color=COLORS["dark"], align="center")
        rect(slide, 3.7, 0.75, 6.1, 0.32, COLORS["teal_dark"])
        tx(slide, 3.7, 0.84, 6.1, 0.12, ss["tableTitle"], size=10, bold=True, color=COLORS["white"], align="center")
        write_table(
            slide,
            3.7,
            1.1,
            [1.5, 0.75, 1.25, 0.85, 0.7, 1.05],
            0.17,
            ss["tableHeaders"],
            ss["ideations"][:18],
            font_size=5.4,
        )
        footer(slide)

    def add_product_slide(prs, item: dict[str, Any]):
        slide = prs.slides.add_slide(blank)
        accent(slide)
        perf = item.get("performance") or {}
        variants = item.get("variants") or []
        gate0 = item.get("gate0") or {}
        comp = item.get("competitor") or {}
        direct = item.get("directCompetitor") or {}
        title = _display_slide_title(item.get("familyName"))
        title_size = 20 if len(title) < 62 else 17
        tx(slide, 0.4, 0.1, 7.2, 0.35, title, size=title_size, bold=True, color=COLORS["dark"])
        tx(slide, 7.8, 0.18, 1.8, 0.16, f"Total SKU Count: {max(1, len(variants))} variations", size=9, color=COLORS["med_gray"], align="right")

        rect(slide, 0.4, 0.6, 2.8, 2.1, COLORS["light_gray"], RGBColor(203, 213, 225))
        tx(slide, 0.58, 1.36, 2.44, 0.18, _clean_text(item.get("comparableSKU") or "Reference TBD", 34), size=8.5, bold=True, color=COLORS["body"], align="center")
        tx(slide, 0.58, 1.58, 2.44, 0.18, "[Product Image]", size=8, color=COLORS["med_gray"], align="center")

        section_header(slide, 3.45, 0.55, 6.2, "Strategy:")
        strategy = gate0.get("overallRead") or gate0.get("sourceEvidence") or (variants[0].get("potentialSKU") if variants else "")
        tx(slide, 3.45, 0.76, 6.2, 0.68, _clean_text(strategy, 610), size=7.2, color=COLORS["body"])

        section_header(slide, 3.45, 1.5, 3.7, "Target Specs:")
        specs = [
            f"Wattage: {item.get('wattages') or 'TBD'}",
            f"CCTs: {item.get('ccts') or 'TBD'}",
            f"Dimming: {item.get('dimming') or 'TBD'}",
            f"Voltage: {item.get('voltage') or 'TBD'}",
            f"CRI: {item.get('cri') or 'TBD'}",
            f"Moisture: {item.get('moistureRating') or 'TBD'}",
            f"Certifications: {item.get('certifications') or 'TBD'}",
        ]
        paragraph_box(slide, 3.45, 1.7, 3.7, 1.0, specs, size=7.0)

        section_header(slide, 7.3, 1.5, 2.35, "Gate 0 Read:", size=9)
        write_table(
            slide,
            7.3,
            1.72,
            [1.2, 1.15],
            0.22,
            ["Metric", "Value"],
            [
                ["Outlook", _clean_text(gate0.get("outlook") or "needs review", 22)],
                ["Confidence", _clean_text(gate0.get("confidence") or "n/a", 22)],
                ["Amazon G2", _clean_text(gate0.get("amazonG2") or "n/a", 22)],
            ],
            font_size=6.4,
        )

        section_header(slide, 0.4, 2.58, 3.0, "Product Variants:", size=10)
        variant_rows = []
        for variant in variants[:4]:
            variant_rows.append([
                _clean_text(variant.get("skuName"), 55),
                _clean_text(variant.get("potentialSKU"), 55),
                _clean_text(", ".join(variant.get("differentiator") or []) or variant.get("features"), 70),
            ])
        write_table(
            slide,
            0.4,
            2.8,
            [3.2, 3.0, 3.0],
            0.2,
            ["SKU", "Description", "Key Differentiators"],
            variant_rows,
            font_size=6.5,
        )

        perf_y = 3.26 + max(1, len(variant_rows)) * 0.2
        section_header(slide, 0.4, perf_y, 5.0, "Comparable SKU Performance / Commercial Read:", size=9)
        write_table(
            slide,
            0.4,
            perf_y + 0.2,
            [1.2, 0.8, 0.8, 0.65, 0.8, 0.8, 1.0, 0.65, 0.65, 0.65],
            0.22,
            ["Comparable SKU", "Shopify Rev", "Amazon Rev", "G2", "Units", "Target MSRP", "Margin $", "Margin %", "Outlook", "Conf."],
            [[
                _clean_text(item.get("comparableSKU"), 24),
                _fmt_money(perf.get("shopifyRevenue")),
                _fmt_money(perf.get("amazonRevenue")),
                _clean_text(gate0.get("amazonG2") or "N/A", 10),
                _fmt_value(perf.get("unitsSold")),
                _fmt_value(perf.get("salePrice")),
                _fmt_money(perf.get("marginDollars")),
                _fmt_value(perf.get("marginPercent")),
                _clean_text(gate0.get("outlook") or "N/A", 14),
                _clean_text(gate0.get("confidence") or "N/A", 12),
            ]],
            header_fill=COLORS["teal_dark"],
            font_size=6.2,
        )

        amazon_y = perf_y + 0.73
        tx(slide, 0.4, amazon_y, 5.8, 0.12, f"Amazon Segment Competitor ({comp.get('segment') or comp.get('channel') or 'Verified competitor'})", size=7.0, bold=True, color=COLORS["dark"])
        write_table(
            slide,
            0.4,
            amazon_y + 0.15,
            [1.2, 1.6, 0.85, 1.25, 1.0, 0.75, 0.75],
            0.155,
            ["ASIN / ID", "Brand", "Price", "Product", "Data Source", "Rating", "Reviews"],
            [[
                _clean_text(comp.get("asin") or comp.get("identifier") or "TBD", 20),
                _clean_text(comp.get("brand") or "TBD", 22),
                _clean_text(comp.get("price") or "N/A", 12),
                _clean_text(comp.get("product") or "Verified Step 3 competitor", 40),
                _clean_text(comp.get("dataSource") or "Step 3", 18),
                _clean_text(comp.get("rating") or "N/A", 8),
                _clean_text(comp.get("reviews") or "N/A", 8),
            ]],
            header_fill=COLORS["navy"],
            font_size=5.8,
        )

        direct_y = amazon_y + 0.49
        tx(slide, 0.4, direct_y, 5.8, 0.12, f"Direct Segment Competitor ({direct.get('segment') or 'Direct / B&M competitor'})", size=7.0, bold=True, color=COLORS["dark"])
        write_table(
            slide,
            0.4,
            direct_y + 0.15,
            [1.6, 1.4, 0.8, 5.4],
            0.155,
            ["Product Code", "Brand", "Price", "Product Name"],
            [[
                _clean_text(direct.get("productCode") or direct.get("identifier") or "TBD", 28),
                _clean_text(direct.get("brand") or "TBD", 22),
                _clean_text(direct.get("price") or "N/A", 12),
                _clean_text(direct.get("productName") or direct.get("product") or gate0.get("whyNotStronger") or "Step 3 details should be reviewed.", 90),
            ]],
            header_fill=COLORS["navy"],
            font_size=5.8,
        )
        source_footnote(slide, config.get("marketSnapshot", {}).get("sourceFootnote"))
        footer(slide)

    def add_decision_slide(prs):
        slide = prs.slides.add_slide(blank)
        accent(slide)
        tx(slide, 0.4, 0.18, 9.2, 0.42, "Gate 0 Decision", size=24, bold=True, color=COLORS["dark"], align="center")
        section_header(slide, 0.65, 1.0, 4.2, "Recommended Decision Logic")
        paragraph_box(
            slide,
            0.65,
            1.28,
            4.1,
            1.15,
            [
                "Move forward: strong demand evidence, clear Sunco line fit, and specific PRD/RFQ next inputs.",
                "Hold: weak or stale demand evidence, unclear competitor match, or unresolved pricing gap.",
                "Needs more data: missing Stackline, pricing, line review, or vendor feasibility evidence.",
            ],
            size=8.2,
        )
        section_header(slide, 5.35, 1.0, 4.0, "Next Step If Approved")
        paragraph_box(
            slide,
            5.35,
            1.28,
            4.0,
            1.15,
            [
                "Create PRD/RFQ package for selected rows.",
                "Coordinate pack-size research with data team before RFQ finalization.",
                "Request vendor quote/sample evidence before any PO-ready work begins.",
            ],
            size=8.2,
        )
        write_table(
            slide,
            0.9,
            3.05,
            [2.2, 2.2, 2.2, 2.2],
            0.28,
            ["Decision", "When To Use", "Owner", "Next Artifact"],
            [
                ["Pursue", "Evidence supports PRD/RFQ", "PM", "PRD / RFQ brief"],
                ["Hold", "Evidence or fit is weak", "PM + Team", "Research gap list"],
                ["More Data", "Need market or line validation", "Data / PM", "Updated Step 3 report"],
            ],
            header_fill=COLORS["teal_dark"],
            font_size=7.2,
        )
        source_footnote(slide)
        footer(slide)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    add_snapshot_slide(prs)
    add_line_review_slide(prs)
    add_summary_slide(prs)
    for item in slide_data:
        add_product_slide(prs, item)
    add_decision_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def generate_gate0_deck(paths: ProjectPaths, category: Category) -> Gate0DeckResult:
    paths.ensure()
    run_timestamp = timestamp()
    warnings: list[str] = []

    gap_path = latest_gap_workbook(paths, category)
    prd_path = latest_prd_ideation_workbook(paths, category)
    research_path = latest_research_report(paths, category)

    if not gap_path:
        warnings.append(f"No Step 1 gap workbook found for {category.run_name}.")
    if not prd_path:
        raise FileNotFoundError(f"No Step 2 PRD ideation workbook found for {category.run_name}. Run Step 2 first.")
    if not research_path:
        warnings.append(f"No Step 3 research report found for {category.run_name}. Deck will use Step 1/2 evidence only.")

    recommendations = _read_table(gap_path, "Recommendations") if gap_path else []
    amazon_recommendations = _read_table(gap_path, "Amazon Recommendations") if gap_path else []
    line_review_rows = _top_line_review_rows(_read_table(gap_path, "Existing SKU Line Review") if gap_path else [])
    if gap_path and not line_review_rows:
        warnings.append("Step 1 workbook has no populated Existing SKU Line Review rows.")

    prd_rows = _read_prd_rows(prd_path)
    if not prd_rows:
        raise ValueError(f"No populated Step 2 ideation rows found in {prd_path}.")

    source_mapping = _read_source_mapping(prd_path)
    research_summary = _read_research_summary(research_path)
    research_details = _read_research_details(research_path)
    ideations = _merge_ideations(
        prd_rows=prd_rows,
        source_mapping=source_mapping,
        research_summary=research_summary,
        research_details=research_details,
    )

    package_dir = paths.cache / "gate0_decks" / category.slug / run_timestamp
    package_dir.mkdir(parents=True, exist_ok=True)
    spec_snapshot_dir = _copy_deck_specs(package_dir)
    output_dir = paths.leadership_deck_category_outputs(category.slug)
    deck_path = output_dir / f"{category.slug}_gate0_deck_{run_timestamp}.pptx"

    paths_used = {
        "step1": str(gap_path) if gap_path else "",
        "step2": str(prd_path),
        "step3": str(research_path) if research_path else "",
    }
    config = _build_config(
        category=category,
        run_timestamp=run_timestamp,
        ideations=ideations,
        line_review_rows=line_review_rows,
        paths_used=paths_used,
    )
    slide_data = _build_slide_data(ideations)

    config_path = package_dir / f"{category.slug}_gate0_config_{run_timestamp}.json"
    slide_data_path = package_dir / f"{category.slug}_gate0_slide_data_{run_timestamp}.json"
    audit_path = package_dir / f"{category.slug}_gate0_audit_{run_timestamp}.json"

    config_path.write_text(json.dumps(config, indent=2), encoding="utf-8")
    slide_data_path.write_text(json.dumps(slide_data, indent=2), encoding="utf-8")
    audit_path.write_text(
        json.dumps(
            {
                "category": category.run_name,
                "generated_at": run_timestamp,
                "source_policy": "Only current Ideation Development Step 1, Step 2, and Step 3 outputs were read.",
                "source_files": paths_used,
                "markdown_spec_source": str(_deck_spec_source_dir()),
                "markdown_spec_snapshot": str(spec_snapshot_dir) if spec_snapshot_dir else "",
                "step1_recommendation_rows": len(recommendations),
                "step1_amazon_recommendation_rows": len(amazon_recommendations),
                "step1_line_review_rows_used": len(line_review_rows),
                "step2_ideation_rows": len(prd_rows),
                "step3_summary_rows_matched": len(research_summary),
                "step3_detail_sheets_matched": len(research_details),
                "warnings": warnings,
            },
            indent=2,
        ),
        encoding="utf-8",
    )

    _write_pptx(
        output_path=deck_path,
        category=category,
        config=config,
        slide_data=slide_data,
        line_review_rows=line_review_rows,
    )

    return Gate0DeckResult(
        deck_path=deck_path,
        package_dir=package_dir,
        config_path=config_path,
        slide_data_path=slide_data_path,
        audit_path=audit_path,
        spec_snapshot_dir=spec_snapshot_dir,
        warnings=warnings,
    )
